import csv
import io
import json
import os
from typing import Any, Callable, Dict, List, Optional

from bs4 import BeautifulSoup
from docx import Document
from fastapi import FastAPI, File, Form, HTTPException, UploadFile
from openpyxl import load_workbook
from pypdf import PdfReader
from pptx import Presentation

app = FastAPI(title="Local Document Parser", version="1.0.0")

MAX_FILE_SIZE_MB = int(os.getenv("PARSER_MAX_FILE_SIZE_MB", "20"))
MAX_FILE_SIZE_BYTES = MAX_FILE_SIZE_MB * 1024 * 1024

SUPPORTED_EXTENSIONS = {
    ".txt",
    ".md",
    ".csv",
    ".json",
    ".pdf",
    ".docx",
    ".xlsx",
    ".pptx",
    ".html",
    ".htm",
}


def _decode_text(content: bytes) -> str:
    for enc in ("utf-8", "utf-8-sig", "big5", "cp950", "latin-1"):
        try:
            return content.decode(enc)
        except UnicodeDecodeError:
            continue
    return content.decode("utf-8", errors="replace")


def _parse_txt(content: bytes) -> str:
    return _decode_text(content).strip()


def _parse_csv(content: bytes) -> str:
    text = _decode_text(content)
    reader = csv.reader(io.StringIO(text))
    lines: List[str] = []
    for row in reader:
        lines.append("\t".join(col.strip() for col in row))
    return "\n".join(lines).strip()


def _parse_json(content: bytes) -> str:
    obj = json.loads(_decode_text(content))
    return json.dumps(obj, ensure_ascii=False, indent=2)


def _parse_pdf(content: bytes) -> str:
    reader = PdfReader(io.BytesIO(content))
    pages: List[str] = []
    for idx, page in enumerate(reader.pages, start=1):
        page_text = page.extract_text() or ""
        pages.append(f"[Page {idx}]\n{page_text.strip()}")
    return "\n\n".join(pages).strip()


def _parse_docx(content: bytes) -> str:
    doc = Document(io.BytesIO(content))
    paragraphs = [p.text.strip() for p in doc.paragraphs if p.text and p.text.strip()]
    return "\n".join(paragraphs).strip()


def _parse_xlsx(content: bytes) -> str:
    wb = load_workbook(io.BytesIO(content), data_only=True)
    output: List[str] = []
    for sheet in wb.worksheets:
        output.append(f"[Sheet] {sheet.title}")
        for row in sheet.iter_rows(values_only=True):
            values = ["" if v is None else str(v).strip() for v in row]
            if any(values):
                output.append("\t".join(values))
        output.append("")
    return "\n".join(output).strip()


def _parse_pptx(content: bytes) -> str:
    prs = Presentation(io.BytesIO(content))
    output: List[str] = []
    for idx, slide in enumerate(prs.slides, start=1):
        output.append(f"[Slide {idx}]")
        for shape in slide.shapes:
            shape_any: Any = shape
            text_value = getattr(shape_any, "text", None)
            if text_value is None:
                text_frame = getattr(shape_any, "text_frame", None)
                if text_frame is not None:
                    text_value = getattr(text_frame, "text", None)

            if isinstance(text_value, str):
                text = text_value.strip()
                if text:
                    output.append(text)
        output.append("")
    return "\n".join(output).strip()


def _parse_html(content: bytes) -> str:
    soup = BeautifulSoup(_decode_text(content), "html.parser")
    return soup.get_text(separator="\n", strip=True)


def _get_parser(filename: str) -> Callable[[bytes], str]:
    ext = os.path.splitext(filename.lower())[1]
    parser_map: Dict[str, Callable[[bytes], str]] = {
        ".txt": _parse_txt,
        ".md": _parse_txt,
        ".csv": _parse_csv,
        ".json": _parse_json,
        ".pdf": _parse_pdf,
        ".docx": _parse_docx,
        ".xlsx": _parse_xlsx,
        ".pptx": _parse_pptx,
        ".html": _parse_html,
        ".htm": _parse_html,
    }
    if ext not in parser_map:
        raise HTTPException(
            status_code=415,
            detail=f"Unsupported file type: {ext}. Supported: {', '.join(sorted(SUPPORTED_EXTENSIONS))}",
        )
    return parser_map[ext]


@app.get("/health")
def health() -> dict:
    return {"ok": True}


@app.get("/supported-types")
def supported_types() -> dict:
    return {"extensions": sorted(SUPPORTED_EXTENSIONS)}


@app.post("/parse")
async def parse_document(
    file: UploadFile = File(...),
    include_metadata: bool = Form(default=True),
    source: Optional[str] = Form(default=None),
) -> dict:
    if not file.filename:
        raise HTTPException(status_code=400, detail="Missing filename")

    content = await file.read()
    if not content:
        raise HTTPException(status_code=400, detail="Empty file")

    if len(content) > MAX_FILE_SIZE_BYTES:
        raise HTTPException(
            status_code=413,
            detail=f"File too large. Max size is {MAX_FILE_SIZE_MB} MB",
        )

    parser = _get_parser(file.filename)
    try:
        text = parser(content)
    except json.JSONDecodeError as exc:
        raise HTTPException(status_code=400, detail=f"Invalid JSON: {exc}") from exc
    except Exception as exc:
        raise HTTPException(status_code=400, detail=f"Failed to parse file: {exc}") from exc

    response = {
        "filename": file.filename,
        "text": text,
        "char_count": len(text),
    }

    if include_metadata:
        response["metadata"] = {
            "content_type": file.content_type,
            "size_bytes": len(content),
            "source": source,
        }

    return response
